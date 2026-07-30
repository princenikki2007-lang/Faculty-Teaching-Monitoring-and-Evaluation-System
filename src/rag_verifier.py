import os
from pathlib import Path
import pdfplumber
from pptx import Presentation

class SyllabusRAGVerifier:
    def __init__(self):
        pass

    def extract_syllabus_text(self, file_path: str) -> str:
        """Extracts plain text content from PDF or PPTX syllabus files."""
        ext = Path(file_path).suffix.lower()
        extracted_text = ""

        try:
            if ext == ".pdf":
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        page_text = page.extract_text()
                        if page_text:
                            extracted_text += page_text + "\n"

            elif ext in [".pptx", ".ppt"]:
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            extracted_text += shape.text + "\n"

            return extracted_text.strip()

        except Exception as e:
            print(f"[RAGVerifier Error] Failed to read {file_path}: {e}")
            return ""